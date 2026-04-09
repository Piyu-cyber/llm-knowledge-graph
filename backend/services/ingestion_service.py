"""
OmniProf v3.0 — Multi-Format Ingestion Service
Handles PDF, DOCX, PPTX, and plain text file ingestion
"""

import os
import json
import logging
from datetime import datetime
from typing import Dict, List, Tuple
from pypdf import PdfReader

from backend.services.jina_multimodal_service import JinaMultimodalService

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


class MultiFormatExtractor:
    """Handles text extraction from multiple file formats"""
    
    # Supported file extensions
    SUPPORTED_FORMATS = {
        '.pdf': 'PDF',
        '.docx': 'DOCX',
        '.doc': 'DOC',
        '.pptx': 'PPTX',
        '.ppt': 'PPT',
        '.txt': 'Text',
        '.png': 'Image',
        '.jpg': 'Image',
        '.jpeg': 'Image',
        '.webp': 'Image',
        '.gif': 'Image'
    }
    
    @staticmethod
    def get_file_format(file_path: str) -> str:
        """Determine file format from extension"""
        _, ext = os.path.splitext(file_path.lower())
        return MultiFormatExtractor.SUPPORTED_FORMATS.get(ext, None)
    
    @staticmethod
    def extract_from_pdf(file_path: str) -> str:
        """Extract text from PDF files"""
        try:
            reader = PdfReader(file_path)
            text_chunks = []
            
            for i, page in enumerate(reader.pages):
                try:
                    extracted = page.extract_text()
                    if extracted:
                        text_chunks.append(extracted)
                except Exception as e:
                    logger.warning(f"⚠️ Skipping page {i}: {str(e)}")
            
            if not text_chunks:
                raise ValueError("No text could be extracted from any page")
            
            full_text = "\n".join(text_chunks)
            return full_text.strip()
        
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
    
    @staticmethod
    def extract_from_docx(file_path: str) -> str:
        """Extract text from DOCX files"""
        if not DOCX_AVAILABLE:
            raise Exception("python-docx not installed. Run: pip install python-docx")
        
        try:
            doc = Document(file_path)
            text_chunks = []
            
            # Extract from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_chunks.append(paragraph.text)
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_chunks.append(" | ".join(row_text))
            
            if not text_chunks:
                raise ValueError("No text could be extracted from document")
            
            full_text = "\n".join(text_chunks)
            return full_text.strip()
        
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
    
    @staticmethod
    def extract_from_pptx(file_path: str) -> str:
        """Extract text from PPTX (PowerPoint) files"""
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx not installed. Run: pip install python-pptx")
        
        try:
            prs = Presentation(file_path)
            text_chunks = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Add slide number as section marker
                text_chunks.append(f"\n--- Slide {slide_num} ---\n")
                
                # Extract from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_chunks.append(shape.text)
                    elif hasattr(shape, "table"):
                        # Handle tables in slides
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_chunks.append(" | ".join(row_text))
            
            if len(text_chunks) <= len(prs.slides):  # Only slide markers, no content
                raise ValueError("No text could be extracted from slides")
            
            full_text = "\n".join(text_chunks)
            return full_text.strip()
        
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    @staticmethod
    def extract_from_txt(file_path: str) -> str:
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read().strip()
            
            if not text:
                raise ValueError("Text file is empty")
            
            return text
        
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    text = f.read().strip()
                return text
            except Exception as e:
                raise Exception(f"Text file extraction failed: {str(e)}")
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, str]:
        """
        Extract text from any supported file format.
        
        Args:
            file_path: Path to the file
        
        Returns:
            Tuple of (text, file_format)
        
        Raises:
            Exception if format not supported or extraction fails
        """
        file_format = MultiFormatExtractor.get_file_format(file_path)
        
        if not file_format:
            raise Exception(
                f"Unsupported file format. Supported: {', '.join(MultiFormatExtractor.SUPPORTED_FORMATS.keys())}"
            )
        
        if file_format == 'PDF':
            text = MultiFormatExtractor.extract_from_pdf(file_path)
        elif file_format in ['DOCX', 'DOC']:
            text = MultiFormatExtractor.extract_from_docx(file_path)
        elif file_format in ['PPTX', 'PPT']:
            text = MultiFormatExtractor.extract_from_pptx(file_path)
        elif file_format == 'Text':
            text = MultiFormatExtractor.extract_from_txt(file_path)
        else:
            raise Exception(f"Unsupported format: {file_format}")
        
        return text, file_format

    @staticmethod
    def extract_content_units(file_path: str, embedding_service: JinaMultimodalService) -> Tuple[List[Dict], str]:
        """
        Convert any supported file into modality-agnostic content units.

        Each content unit includes source reference, page/slide index, and modality tag.
        """
        file_format = MultiFormatExtractor.get_file_format(file_path)
        if not file_format:
            raise Exception(
                f"Unsupported file format. Supported: {', '.join(MultiFormatExtractor.SUPPORTED_FORMATS.keys())}"
            )

        source = os.path.basename(file_path)
        units: List[Dict] = []

        if file_format == "PDF":
            reader = PdfReader(file_path)
            for idx, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    units.append(
                        {
                            "source_ref": source,
                            "position": idx,
                            "modality": "text",
                            "content": text,
                            "embedding": embedding_service.embed_text(text),
                        }
                    )

        elif file_format in ["DOCX", "DOC"]:
            if not DOCX_AVAILABLE:
                raise Exception("python-docx not installed. Run: pip install python-docx")
            doc = Document(file_path)
            for idx, paragraph in enumerate(doc.paragraphs, start=1):
                text = paragraph.text.strip()
                if text:
                    units.append(
                        {
                            "source_ref": source,
                            "position": idx,
                            "modality": "text",
                            "content": text,
                            "embedding": embedding_service.embed_text(text),
                        }
                    )

        elif file_format in ["PPTX", "PPT"]:
            if not PPTX_AVAILABLE:
                raise Exception("python-pptx not installed. Run: pip install python-pptx")
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = (shape.text or "").strip()
                        if text:
                            units.append(
                                {
                                    "source_ref": source,
                                    "position": slide_idx,
                                    "modality": "text",
                                    "content": text,
                                    "embedding": embedding_service.embed_text(text),
                                }
                            )
                    if getattr(shape, "shape_type", None) is not None and hasattr(shape, "image"):
                        # Diagram/image unit, no OCR text extraction.
                        units.append(
                            {
                                "source_ref": source,
                                "position": slide_idx,
                                "modality": "image",
                                "content": "",
                                "embedding": embedding_service.embed_diagram(file_path, description=f"slide_{slide_idx}_diagram"),
                            }
                        )

        elif file_format == "Image":
            units.append(
                {
                    "source_ref": source,
                    "position": 1,
                    "modality": "image",
                    "content": "",
                    "embedding": embedding_service.embed_image(file_path),
                }
            )

        else:  # Text
            text = MultiFormatExtractor.extract_from_txt(file_path)
            units.append(
                {
                    "source_ref": source,
                    "position": 1,
                    "modality": "text",
                    "content": text,
                    "embedding": embedding_service.embed_text(text),
                }
            )

        return units, file_format


class IngestionService:
    """Handles document ingestion and knowledge extraction"""
    
    def __init__(self, llm_service, rag_service, graph_service):
        self.llm_service = llm_service
        self.rag_service = rag_service
        self.graph_service = graph_service
        self.extractor = MultiFormatExtractor()
        self.embedding_service = JinaMultimodalService()
        self.review_queue_path = os.path.join("data", "review_queue.json")
        # Enable strict pre-write validation for Phase 2 acceptance/integration flows.
        self.enable_prewrite_validation = False
    
    def _normalize_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        # Remove excessive whitespace
        text = "\n".join(line.strip() for line in text.split("\n") if line.strip())
        # Limit length for LLM processing
        MAX_CHARS = 15000
        return text[:MAX_CHARS]
    
    def _clean_text(self, value):
        """Clean text values"""
        if not value:
            return ""
        
        value = value.strip()
        
        # Remove common prefixes
        if value.lower().startswith("name:"):
            value = value.split(":", 1)[1].strip()
        
        return value

    def _append_review_queue(self, source_doc: str, errors: List[Dict]) -> None:
        os.makedirs(os.path.dirname(self.review_queue_path), exist_ok=True)
        payload = []
        if os.path.exists(self.review_queue_path):
            try:
                with open(self.review_queue_path, "r", encoding="utf-8") as f:
                    payload = json.load(f)
            except Exception:
                payload = []
        payload.append(
            {
                "timestamp": datetime.now().isoformat(),
                "source_doc": source_doc,
                "errors": errors,
            }
        )
        with open(self.review_queue_path, "w", encoding="utf-8") as f:
            json.dump(payload, f, indent=2)

    def _prevalidate_extraction(self, llm_data: Dict, source_doc: str) -> List[Dict]:
        """Validate extracted graph payload before canonical write."""
        issues: List[Dict] = []
        nodes = llm_data.get("nodes", [])
        edges = llm_data.get("edges", [])

        # Duplicate concept names in extraction payload.
        seen_concepts = set()
        for node in nodes:
            if node.get("level", "").upper() != "CONCEPT":
                continue
            key = node.get("name", "").strip().lower()
            if key in seen_concepts:
                issues.append({"type": "duplicate_concept", "name": node.get("name", "")})
            seen_concepts.add(key)

        # Cycle detection on REQUIRES edges in payload.
        adj: Dict[str, List[str]] = {}
        for edge in edges:
            if edge.get("type", "").upper() != "REQUIRES":
                continue
            source = edge.get("source", "")
            target = edge.get("target", "")
            if source and target:
                adj.setdefault(source, []).append(target)

        visited = set()
        stack = set()

        def dfs(node: str) -> None:
            visited.add(node)
            stack.add(node)
            for nxt in adj.get(node, []):
                if nxt not in visited:
                    dfs(nxt)
                elif nxt in stack:
                    issues.append({"type": "prerequisite_cycle", "at": node, "next": nxt})
            stack.remove(node)

        for n in list(adj.keys()):
            if n not in visited:
                dfs(n)

        # Orphan detection in payload graph for non-module nodes.
        names = [n.get("name", "") for n in nodes if n.get("name")]
        connected = set()
        for edge in edges:
            connected.add(edge.get("source", ""))
            connected.add(edge.get("target", ""))
        for node in nodes:
            if node.get("level", "").upper() == "MODULE":
                continue
            name = node.get("name", "")
            if name and name not in connected:
                issues.append({"type": "orphan_node", "name": name, "level": node.get("level")})

        if issues:
            self._append_review_queue(source_doc, issues)
        return issues

    def normalize_to_content_units(self, file_path: str) -> Dict:
        """Public Phase 2 API: normalize source file into modality-agnostic units."""
        units, file_format = self.extractor.extract_content_units(file_path, self.embedding_service)
        return {
            "status": "success",
            "file_format": file_format,
            "content_units": units,
            "unit_count": len(units),
        }

    def extract_graph_from_units(self, content_units: List[Dict], course_guide: str = "") -> Dict:
        """Extract graph payload with one LLM call per content unit."""
        merged_nodes: List[Dict] = []
        merged_edges: List[Dict] = []
        calls = 0

        for unit in content_units:
            modality = unit.get("modality", "text")
            if modality != "text":
                continue
            text = (unit.get("content") or "").strip()
            if not text:
                continue
            calls += 1
            payload = self.llm_service.extract_concepts_hierarchical(text, guide=course_guide)
            merged_nodes.extend(payload.get("nodes", []))
            merged_edges.extend(payload.get("edges", []))

        return {"nodes": merged_nodes, "edges": merged_edges, "llm_calls": calls}
    
    def ingest(self, file_path: str, course_owner: str = "system", course_guide_text: str = "", source_doc_name: str = "") -> Dict:
        """
        Ingest a document in any supported format.
        
        Args:
            file_path: Path to the document
            course_owner: Course owner ID for hierarchy
        
        Returns:
            Dict with ingestion results including validation errors
        """
        try:
            # Step 1: Reset RAG
            self.rag_service.reset()
            
            # Step 2: Build modality-agnostic content units
            logger.info(f"Normalizing content from {file_path}")
            normalized = self.normalize_to_content_units(file_path)
            file_format = normalized["file_format"]
            units = normalized["content_units"]
            source_doc_name = source_doc_name or os.path.basename(file_path)
            text_units = [u.get("content", "") for u in units if u.get("modality") == "text" and u.get("content")]
            text = "\n".join(text_units)
            
            if not text:
                # Image-only docs are valid for multimodal indexing and should not crash ingestion.
                text = ""
            
            # Step 3: Normalize text for LLM extraction
            text = self._normalize_text(text) if text else ""
            
            # Step 4: Extract hierarchical concepts using one call per text content unit
            logger.info("Extracting hierarchical concepts using LLM")
            llm_data = self.extract_graph_from_units(units, course_guide=course_guide_text)
            
            if not llm_data or not llm_data.get("nodes"):
                return {
                    "status": "error",
                    "message": "Failed to extract concepts from document"
                }

            # Step 4b: Pre-write validation (cycles, orphans, duplicates)
            if self.enable_prewrite_validation:
                prewrite_issues = self._prevalidate_extraction(llm_data, os.path.basename(file_path))
                if prewrite_issues:
                    return {
                        "status": "error",
                        "message": "Extraction failed validation",
                        "validation_errors": prewrite_issues,
                        "review_queued": True,
                    }
            
            # Step 5: Insert into graph with hierarchical structure
            logger.info("Inserting into knowledge graph")
            insert_result = self.graph_service.insert_from_llm_hierarchical(
                data=llm_data,
                course_owner=course_owner,
                source_doc=source_doc_name,
                file_format=file_format
            )
            
            if insert_result["status"] != "success":
                return insert_result
            
            # Step 6: Validate graph integrity
            logger.info("Validating graph integrity")
            validation_result = self.graph_service.validate_graph()
            
            # Step 7: Store full text in RAG
            if text:
                try:
                    self.rag_service.ingest_documents(text, guide_text=course_guide_text, source_name=source_doc_name)
                    logger.info("Text stored in RAG system")
                except Exception as e:
                    logger.warning(f"⚠️ RAG ingestion failed: {str(e)}")
            
            # Return results with validation info
            return {
                "status": "success",
                "file_format": file_format,
                "concepts_added": insert_result.get("concepts_added", 0),
                "relationships_added": insert_result.get("relationships_added", 0),
                "facts_added": insert_result.get("facts_added", 0),
                "modules_added": insert_result.get("modules_added", 0),
                "topics_added": insert_result.get("topics_added", 0),
                "validation": {
                    "is_valid": validation_result["status"] == "valid",
                    "issue_count": validation_result.get("issue_count", 0),
                    "issues": validation_result.get("issues", [])
                },
                "content_units": len(units),
                "llm_calls": llm_data.get("llm_calls", 0),
            }
        
        except Exception as e:
            logger.error(f"Ingestion error: {str(e)}")
            return {
                "status": "error",
                "message": str(e)
            }

    def ingest_incremental(self, file_path: str, course_owner: str = "system", course_guide_text: str = "", source_doc_name: str = "") -> Dict:
        """Incrementally re-ingest one source document while preserving unrelated overlays."""
        try:
            normalized = self.normalize_to_content_units(file_path)
            file_format = normalized["file_format"]
            units = normalized["content_units"]
            source_doc = source_doc_name or os.path.basename(file_path)
            llm_data = self.extract_graph_from_units(units, course_guide=course_guide_text)

            if self.enable_prewrite_validation:
                prewrite_issues = self._prevalidate_extraction(llm_data, source_doc)
                if prewrite_issues:
                    return {
                        "status": "error",
                        "message": "Extraction failed validation",
                        "validation_errors": prewrite_issues,
                        "review_queued": True,
                    }

            result = self.graph_service.incremental_reingest_from_llm(
                data=llm_data,
                course_owner=course_owner,
                source_doc=source_doc,
                file_format=file_format,
            )

            result["content_units"] = len(units)
            result["llm_calls"] = llm_data.get("llm_calls", 0)
            return result
        except Exception as e:
            logger.error(f"Incremental ingestion error: {str(e)}")
            return {"status": "error", "message": str(e)}
